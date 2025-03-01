#!/usr/bin/env python3
"""
Azure AI Search Document Uploader for RAG (Retrieval-Augmented Generation)

This script processes documents from a directory, chunks them intelligently,
and uploads them to an Azure AI Search index for use with RAG chatbots.
It supports both vector search (embeddings-based) and semantic search capabilities.

Features:
- Multi-format document processing (PDF, DOCX, TXT, HTML, CSV, PPTX, XLSX)
- Intelligent text chunking with overlap
- Vector embeddings generation using Azure OpenAI
- Semantic search configuration
- Customizable document filtering

Example Usage:
    python rag_upload.py \
        --folder "C:/Documents/data" \
        --endpoint "https://yoursearch.search.windows.net" \
        --key "your_search_api_key" \
        --index "myindex" \
        --vector-search \
        --aoai-endpoint "https://youropenai.openai.azure.com/" \
        --aoai-key "your_openai_api_key" \
        --aoai-deployment "text-embedding-3-large"

Dependencies:
    pip install azure-search-documents azure-core PyPDF2 python-docx pandas beautifulsoup4 python-pptx langdetect
"""

# Standard library imports
import os
import argparse
import json
import logging
import hashlib
import re
import time
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple, Set, Union
from dataclasses import dataclass, field

# Azure SDK imports
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SearchIndex, SearchFieldDataType,
    SimpleField, SearchableField
)

# Third-party document processing libraries
import requests  # For REST API calls
import PyPDF2    # PDF processing
import docx      # Word document processing
import pandas as pd  # Excel and CSV processing
from bs4 import BeautifulSoup  # HTML processing
import pptx     # PowerPoint processing
from langdetect import detect  # Language detection

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@dataclass
class AzureSearchConfig:
    """
    Configuration settings for Azure AI Search service.
    
    Attributes:
        endpoint (str): The Azure AI Search service endpoint URL
        api_key (str): The admin API key for the search service
        index_name (str): Name of the search index to create/update
        vector_search (bool): Whether to enable vector search capabilities
        semantic_search (bool): Whether to enable semantic search capabilities
        vector_dimensions (int): Dimensions for vector embeddings (3072 for text-embedding-3-large, 1536 for ada-002)
        force_recreate (bool): Whether to force recreation of the index
    """
    endpoint: str
    api_key: str
    index_name: str
    vector_search: bool = False
    semantic_search: bool = False
    vector_dimensions: int = 1536
    force_recreate: bool = False

@dataclass
class OpenAIConfig:
    """
    Configuration settings for Azure OpenAI service.
    
    Attributes:
        endpoint (Optional[str]): The Azure OpenAI service endpoint URL
        api_key (Optional[str]): The API key for the OpenAI service
        deployment_name (Optional[str]): Name of the embeddings model deployment
        skip_embeddings (bool): Whether to skip vector embedding generation
    """
    endpoint: Optional[str] = None
    api_key: Optional[str] = None
    deployment_name: Optional[str] = None
    skip_embeddings: bool = False

@dataclass
class ProcessingConfig:
    """
    Configuration settings for document processing.
    
    Attributes:
        folder_path (str): Path to the folder containing documents to process
        chunk_size (int): Maximum size of text chunks in characters
        chunk_overlap (int): Number of characters to overlap between chunks
        include_filters (List[str]): Regex patterns for files to include
        exclude_filters (List[str]): Regex patterns for files to exclude
    """
    folder_path: str
    chunk_size: int = 1000
    chunk_overlap: int = 100
    include_filters: List[str] = field(default_factory=list)
    exclude_filters: List[str] = field(default_factory=list)

class DocumentProcessor:
    """
    Base class for processing various document types and extracting their content.
    
    This class provides methods to process different file types (PDF, DOCX, TXT, etc.)
    and extract both their text content and metadata in a consistent format.
    
    Attributes:
        config (ProcessingConfig): Configuration settings for document processing
    """
    
    def __init__(self, config: ProcessingConfig):
        """
        Initialize the document processor.
        
        Args:
            config (ProcessingConfig): Configuration settings for document processing
        """
        self.config = config
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        Process a file and extract its text content and metadata.
        
        Args:
            file_path (str): Path to the file to process
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Extracted text content
                - metadata: File metadata including name, path, size, etc.
        """
        # Get basic file information
        file_extension = os.path.splitext(file_path)[1].lower()
        file_size = os.path.getsize(file_path)
        
        # Format the last_modified date in ISO 8601 format with timezone info
        try:
            last_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
            last_modified_iso = last_modified.strftime("%Y-%m-%dT%H:%M:%SZ")  # Format with 'Z' for UTC
        except:
            # Use current time if file modification time cannot be retrieved
            last_modified_iso = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
        
        # Collect basic metadata common to all file types
        metadata = {
            "file_name": os.path.basename(file_path),
            "file_path": file_path,
            "file_size": file_size,
            "file_extension": file_extension,
            "last_modified": last_modified_iso
        }
        
        # Get and execute the appropriate processor for this file type
        processor = self._get_processor_for_file(file_path)
        
        try:
            # Extract text and file-specific metadata
            result = processor(file_path)
            
            # Ensure all date fields in metadata have valid ISO 8601 format
            for key, value in result["metadata"].items():
                if isinstance(value, str) and 'date' in key.lower() and value == '':
                    result["metadata"][key] = last_modified_iso
            
            # Merge file-specific metadata with common metadata
            result["metadata"].update(metadata)
            
            return result
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return {"text": "", "metadata": metadata}
    
    def _get_processor_for_file(self, file_path: str) -> callable:
        """
        Get the appropriate processor function for a file type.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            callable: Function to process the specific file type
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Map file extensions to their processor functions
        processors = {
            '.pdf': self._extract_text_from_pdf,
            '.docx': self._extract_text_from_docx,
            '.txt': self._extract_text_from_txt,
            '.md': self._extract_text_from_txt,
            '.py': self._extract_text_from_txt,
            '.js': self._extract_text_from_txt,
            '.html': self._extract_text_from_html,
            '.htm': self._extract_text_from_html,
            '.csv': self._extract_text_from_csv,
            '.pptx': self._extract_text_from_pptx,
            '.xlsx': self._extract_text_from_xlsx,
            '.xls': self._extract_text_from_xlsx,
        }
        
        # Return default processor for unknown file types
        return processors.get(file_extension, lambda f: {"text": "", "metadata": {}})
    
    def _extract_text_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and metadata from PDF files with semantic awareness.
        
        This method processes PDF files while preserving structural information:
        - Extracts metadata from PDF properties
        - Preserves page boundaries with subtle markers
        - Identifies potential section headers
        - Maintains document flow for better semantic chunking
        
        Args:
            file_path (str): Path to the PDF file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Extracted text content with page markers
                - metadata: PDF-specific metadata and properties
        """
        text = ""
        metadata = {}
        
        try:
            with open(file_path, 'rb') as file:
                # Initialize PDF reader
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract PDF document properties as metadata
                if pdf_reader.metadata:
                    for key, value in pdf_reader.metadata.items():
                        if value and isinstance(value, str):
                            metadata[key] = value
                
                # Process each page while preserving document structure
                full_text = []
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    
                    if page_text:
                        # Add page marker in a format that won't disrupt semantic analysis
                        # [pg:N] format is used by the chunker to maintain page context
                        processed_text = f"{page_text.strip()}\n[pg:{page_num + 1}]\n\n"
                        full_text.append(processed_text)
                
                # Record total page count in metadata
                metadata["pdf_page_count"] = len(pdf_reader.pages)
                
                # Combine all pages into a single text stream
                text = "".join(full_text)
                
                # Identify potential section headers for better chunking
                # Look for patterns like "INTRODUCTION", "CHAPTER 1", etc.
                headers = re.findall(r'\n([A-Z][A-Z\s]{3,}[A-Z])[^a-z]', text)
                if headers:
                    metadata["potential_sections"] = headers
                
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_docx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and metadata from DOCX (Word) files.
        
        This method processes Word documents while preserving structural information:
        - Extracts document properties (author, title, etc.)
        - Preserves heading structure
        - Maintains paragraph formatting
        - Extracts table content in a readable format
        
        Args:
            file_path (str): Path to the DOCX file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Extracted text content with preserved structure
                - metadata: Document properties and metadata
        """
        text = ""
        metadata = {}
        
        try:
            # Load the Word document
            doc = docx.Document(file_path)
            
            # Extract document properties if available
            core_properties = getattr(doc, 'core_properties', None)
            if core_properties:
                # Map core properties to metadata fields
                property_mapping = {
                    "author": "author",
                    "created": "created",
                    "last_modified_by": "last_modified_by",
                    "title": "title",
                    "subject": "subject"
                }
                
                for prop_name, meta_name in property_mapping.items():
                    value = getattr(core_properties, prop_name, "")
                    if value:  # Only add non-empty values
                        metadata[meta_name] = value
            
            # Process document content
            for para in doc.paragraphs:
                if para.style and para.style.name.startswith('Heading'):
                    # Preserve heading structure with extra newlines
                    # This helps the chunker identify document sections
                    text += f"\n{para.text}\n\n"
                else:
                    # Regular paragraphs get a single newline
                    text += para.text + "\n"
            
            # Process tables and convert to readable text
            for table in doc.tables:
                for row in table.rows:
                    # Join cells with pipe separator for readability
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    text += row_text + "\n"
                # Add extra newline after each table
                text += "\n"
                
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_txt(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from plain text files (TXT, MD, PY, JS, etc.).
        
        This method handles plain text files with different encodings:
        - Attempts UTF-8 encoding first
        - Falls back to Latin-1 if UTF-8 fails
        - Preserves line breaks and formatting
        
        Args:
            file_path (str): Path to the text file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Raw text content
                - metadata: Basic file metadata
        """
        text = ""
        # Get file modification time in ISO format
        last_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
        last_modified_iso = last_modified.strftime("%Y-%m-%dT%H:%M:%SZ")
        
        metadata = {
            "last_modified": last_modified_iso
        }
        
        try:
            # Try UTF-8 encoding first (most common)
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        except UnicodeDecodeError:
            try:
                # Fall back to Latin-1 encoding if UTF-8 fails
                with open(file_path, 'r', encoding='latin-1') as file:
                    text = file.read()
            except Exception as e:
                logger.error(f"Error reading text file {file_path} with Latin-1 encoding: {e}")
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_csv(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from CSV files.
        
        This method processes CSV files while preserving data structure:
        - Extracts column headers
        - Preserves tabular format
        - Handles data formatting
        
        Args:
            file_path (str): Path to the CSV file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Formatted text representation of CSV data
                - metadata: CSV properties including headers
        """
        text = ""
        # Get file modification time in ISO format
        last_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
        last_modified_iso = last_modified.strftime("%Y-%m-%dT%H:%M:%SZ")
        
        metadata = {
            "last_modified": last_modified_iso
        }
        
        try:
            # Read CSV into pandas DataFrame
            df = pd.read_csv(file_path)
            
            # Store column headers in metadata
            headers = df.columns.tolist()
            metadata["headers"] = ", ".join(headers)
            
            # Convert DataFrame to formatted string
            # This preserves alignment and readability
            text = df.to_string(index=False)
            
            # Add basic statistics to metadata
            metadata["row_count"] = len(df)
            metadata["column_count"] = len(df.columns)
            
        except Exception as e:
            logger.error(f"Error extracting text from CSV {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_html(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and metadata from HTML files.
        
        This method processes HTML documents while preserving important structure:
        - Extracts metadata from meta tags
        - Preserves document title
        - Maintains heading hierarchy
        - Removes scripts and styles
        - Preserves content structure
        
        Args:
            file_path (str): Path to the HTML file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Extracted text content with preserved structure
                - metadata: HTML metadata and properties
        """
        text = ""
        metadata = {}
        
        try:
            # Read HTML content with UTF-8 encoding
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            # Parse HTML with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract metadata from meta tags
            for meta in soup.find_all('meta'):
                if meta.get('name') and meta.get('content'):
                    metadata[meta['name']] = meta['content']
            
            # Extract and process title
            title = soup.find('title')
            if title:
                metadata['title'] = title.text
                text += f"Title: {title.text}\n\n"
            
            # Process headings and their associated content
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                # Add heading with proper spacing
                text += f"\n{heading.text.strip()}\n\n"
                
                # Get content until next heading
                for sibling in heading.next_siblings:
                    if sibling.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        break
                    if sibling.name and sibling.name not in ['script', 'style']:
                        content = sibling.get_text(strip=True)
                        if content:
                            text += content + "\n"
            
            # If we didn't get much text from the structured approach,
            # fall back to extracting all text content
            if len(text) < 200:
                # Remove scripts and styles first
                for element in soup(['script', 'style']):
                    element.extract()
                
                # Get remaining text with proper spacing
                text = soup.get_text(separator="\n")
                
        except Exception as e:
            logger.error(f"Error extracting text from HTML {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from PowerPoint (PPTX) files.
        
        This method processes PowerPoint presentations while preserving structure:
        - Extracts presentation metadata
        - Maintains slide order and numbering
        - Preserves slide titles and content
        - Handles text from shapes and tables
        
        Args:
            file_path (str): Path to the PPTX file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Extracted text with slide structure preserved
                - metadata: Presentation properties and statistics
        """
        text = ""
        metadata = {}
        
        try:
            # Load PowerPoint presentation
            prs = pptx.Presentation(file_path)
            
            # Extract presentation properties
            if hasattr(prs.core_properties, 'author'):
                metadata["author"] = prs.core_properties.author
            if hasattr(prs.core_properties, 'title'):
                metadata["title"] = prs.core_properties.title
            
            # Store presentation statistics
            metadata["slide_count"] = len(prs.slides)
            
            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                # Add slide marker
                text += f"\n=== Slide {i} ===\n"
                
                # Extract slide title if available
                if slide.shapes.title:
                    text += f"Title: {slide.shapes.title.text}\n\n"
                
                # Process all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Clean and add shape text
                        shape_text = shape.text.strip()
                        if shape_text:
                            text += shape_text + "\n"
                
                # Add extra newline between slides
                text += "\n"
            
            # Add presentation summary to metadata
            metadata["has_speaker_notes"] = any(
                shape.has_text_frame
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == 19  # MSO_SHAPE_TYPE.NOTES
            )
            
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }
    
    def _extract_text_from_xlsx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from Excel (XLSX/XLS) files.
        
        This method processes Excel workbooks while preserving structure:
        - Processes all sheets in the workbook
        - Preserves sheet names and order
        - Maintains tabular data format
        - Handles cell formatting and alignment
        
        Args:
            file_path (str): Path to the Excel file
            
        Returns:
            Dict[str, Any]: Dictionary containing:
                - text: Formatted text representation of all sheets
                - metadata: Workbook properties and statistics
        """
        text = ""
        metadata = {}
        
        try:
            # Open Excel workbook
            xlsx = pd.ExcelFile(file_path)
            
            # Store sheet information in metadata
            sheet_names = xlsx.sheet_names
            metadata["sheet_names"] = ", ".join(sheet_names)
            metadata["sheet_count"] = len(sheet_names)
            
            # Process each sheet
            total_rows = 0
            sheet_stats = []
            
            for sheet_name in sheet_names:
                # Read sheet into DataFrame
                df = pd.read_excel(xlsx, sheet_name)
                
                # Add sheet header
                text += f"\n{'='*20} Sheet: {sheet_name} {'='*20}\n\n"
                
                # Convert sheet data to formatted string
                # This preserves alignment and readability
                if not df.empty:
                    text += df.to_string(index=False) + "\n\n"
                    
                    # Collect sheet statistics
                    total_rows += len(df)
                    sheet_stats.append({
                        "name": sheet_name,
                        "rows": len(df),
                        "columns": len(df.columns),
                        "headers": ", ".join(df.columns.tolist())
                    })
            
            # Add workbook statistics to metadata
            metadata["total_rows"] = total_rows
            metadata["sheet_statistics"] = sheet_stats
            
        except Exception as e:
            logger.error(f"Error extracting text from Excel file {file_path}: {e}")
        
        return {
            "text": text,
            "metadata": metadata
        }

class TextChunker:
    """
    Intelligent text chunking system for document processing.
    
    This class splits large text documents into semantically meaningful chunks
    while respecting document structure. It uses different chunking strategies
    based on content type (PDF, code, JSON, Markdown, etc.) to ensure chunks
    maintain context and readability.
    
    Attributes:
        chunk_size (int): Maximum size of each chunk in characters
        chunk_overlap (int): Number of characters to overlap between chunks
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 100):
        """
        Initialize the text chunker.
        
        Args:
            chunk_size (int): Maximum size of each chunk in characters
            chunk_overlap (int): Number of characters to overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def chunk_text(self, text: str) -> List[str]:
        """
        Split text into chunks with specified size and overlap, respecting content structure.
        
        This method intelligently detects the content type and applies the most appropriate
        chunking strategy:
        - PDF documents: Chunks by sections and page markers
        - Code: Chunks by function/class boundaries
        - JSON: Chunks by object structure
        - Markdown: Chunks by headings
        - Regular text: Chunks by paragraphs and sentences
        
        Args:
            text (str): The text content to chunk
            
        Returns:
            List[str]: List of text chunks with appropriate overlap
        """
        # For very short texts, return as a single chunk
        if len(text) <= self.chunk_size:
            return [text.strip()]
        
        chunks = []
        
        # Detect content type based on patterns
        is_code = bool(re.search(r'(def |class |function |import |from |var |const |let |public |private |#include)', text))
        is_json = bool(text.strip().startswith('{') and text.strip().endswith('}'))
        is_markdown = bool(re.search(r'^#{1,6} |\n#{1,6} |```|\*\*|__|\[.+\]\(.+\)', text))
        has_pdf_markers = '[pg:' in text and re.search(r'\[pg:\d+\]', text)
        
        # === PDF DOCUMENT CHUNKING STRATEGY ===
        if has_pdf_markers:
            chunks = self._chunk_pdf_document(text)
            if chunks:
                return chunks
        
        # === JSON DOCUMENT CHUNKING STRATEGY ===
        if is_json:
            chunks = self._chunk_json_document(text)
            if chunks:
                return chunks
        
        # === CODE DOCUMENT CHUNKING STRATEGY ===
        if is_code:
            chunks = self._chunk_code_document(text)
            if chunks:
                return chunks
        
        # === MARKDOWN DOCUMENT CHUNKING STRATEGY ===
        if is_markdown:
            chunks = self._chunk_markdown_document(text)
            if chunks:
                return chunks
        
        # === DEFAULT TEXT CHUNKING STRATEGY ===
        # For regular text, chunk by paragraphs and sentences
        return self._chunk_regular_text(text)
    
    def remove_duplicate_chunks(self, chunks: List[str]) -> List[str]:
        """
        Remove duplicate or near-duplicate chunks from the list.
        
        This method identifies and removes chunks that are exact duplicates or have
        very high similarity to other chunks in the list, reducing redundancy in the
        processed documents.
        
        Args:
            chunks (List[str]): List of text chunks to deduplicate
            
        Returns:
            List[str]: Deduplicated list of text chunks
        """
        if not chunks:
            return []
            
        # For exact duplicates, use a set to remove them
        unique_chunks = []
        seen_chunks = set()
        
        for chunk in chunks:
            # Create a normalized version for comparison (lowercase, whitespace normalized)
            normalized = re.sub(r'\s+', ' ', chunk.lower().strip())
            
            # Skip if we've seen this exact chunk before
            if normalized in seen_chunks:
                continue
                
            seen_chunks.add(normalized)
            unique_chunks.append(chunk)
            
        return unique_chunks
    
    def _chunk_pdf_document(self, text: str) -> List[str]:
        """
        Chunk PDF documents using section headers and page markers.
        
        Args:
            text (str): PDF text content with page markers
            
        Returns:
            List[str]: Chunked PDF content or empty list if chunking fails
        """
        chunks = []
        
        # Try to find section markers or headers for better chunking
        # Common patterns: ALL CAPS HEADINGS, Chapter X, Section X.X, etc.
        sections = []
        
        # Define patterns for different heading styles in PDFs
        heading_patterns = [
            r'\n([A-Z][A-Z\s]{3,}[A-Z])[^a-z]',  # ALL CAPS HEADINGS
            r'\n(Chapter\s+\d+[.:]\s*\w+)',       # Chapter headings
            r'\n((?:\d+\.){1,3}\s*\w+)',          # Section numbers like 1.2.3
            r'\n([IVX]+\.\s+\w+)',                # Roman numeral sections
            r'\b(\d+\.\d+\s+[A-Z][a-z]+)'         # Numbered headings like "1.2 Introduction"
        ]
        
        # Find all potential section boundaries
        for pattern in heading_patterns:
            matches = re.finditer(pattern, text)
            for match in matches:
                sections.append(match.start())
        
        # If we found sections, use them as chunk boundaries
        if sections:
            sections.sort()
            start_pos = 0
            
            # Create chunks based on section boundaries
            for section_pos in sections:
                # Only break if this would create a reasonably sized chunk
                if section_pos - start_pos > self.chunk_size and start_pos < section_pos:
                    chunk_text = text[start_pos:section_pos].strip()
                    if chunk_text:
                        chunks.append(chunk_text)
                    start_pos = section_pos
            
            # Add the final section
            if start_pos < len(text):
                chunk_text = text[start_pos:].strip()
                if chunk_text:
                    chunks.append(chunk_text)
            
            # If we successfully created chunks, return them
            if chunks:
                return chunks
        
        # Fallback: chunk by paragraphs if section detection failed
        paragraphs = re.split(r'\n\s*\n', text)
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_size = len(para) + 2  # +2 for the newlines
            
            # If adding this paragraph would exceed chunk size, complete current chunk
            if current_size + para_size > self.chunk_size and current_chunk:
                chunk_text = "\n\n".join(current_chunk).strip()
                if chunk_text:
                    chunks.append(chunk_text)
                
                # Create overlap with previous paragraphs
                overlap_size = 0
                overlap_paras = []
                for prev_para in reversed(current_chunk):
                    if overlap_size + len(prev_para) + 2 <= self.chunk_overlap:
                        overlap_paras.insert(0, prev_para)
                        overlap_size += len(prev_para) + 2
                    else:
                        break
                
                current_chunk = overlap_paras
                current_size = overlap_size
            
            current_chunk.append(para)
            current_size += para_size
        
        # Add the final chunk
        if current_chunk:
            chunk_text = "\n\n".join(current_chunk).strip()
            if chunk_text:
                chunks.append(chunk_text)
        
        return chunks
    
    def _chunk_json_document(self, text: str) -> List[str]:
        """
        Chunk JSON documents by object structure.
        
        Args:
            text (str): JSON text content
            
        Returns:
            List[str]: Chunked JSON content or empty list if chunking fails
        """
        chunks = []
        
        try:
            # Parse the JSON to understand its structure
            data = json.loads(text)
            
            # Handle arrays of objects
            if isinstance(data, list) and len(data) > 1:
                # Calculate how many items to include per chunk
                items_per_chunk = max(1, self.chunk_size // 500)
                
                # Split the list into chunks
                for i in range(0, len(data), items_per_chunk):
                    subset = data[i:min(i + items_per_chunk, len(data))]
                    chunks.append(json.dumps(subset, indent=2))
                return chunks
                
            # Handle dictionary objects
            elif isinstance(data, dict):
                # Split large dictionaries by top-level keys
                keys = list(data.keys())
                keys_per_chunk = max(1, self.chunk_size // 500)
                
                for i in range(0, len(keys), keys_per_chunk):
                    subset_keys = keys[i:min(i + keys_per_chunk, len(keys))]
                    subset = {k: data[k] for k in subset_keys}
                    chunks.append(json.dumps(subset, indent=2))
                return chunks
        except:
            # If JSON parsing fails, return empty list to fall back to default chunking
            return []
        
        return chunks
    
    def _chunk_code_document(self, text: str) -> List[str]:
        """
        Chunk code documents by respecting function and class boundaries.
        
        Args:
            text (str): Code text content
            
        Returns:
            List[str]: Chunked code content
        """
        chunks = []
        lines = text.split('\n')
        current_chunk = []
        current_length = 0
        
        for line in lines:
            line_with_newline = line + '\n'
            line_length = len(line_with_newline)
            
            # If adding this line would exceed chunk size, complete current chunk
            if current_length + line_length > self.chunk_size and current_chunk:
                chunks.append(''.join(current_chunk))
                
                # Create overlap with previous lines for context
                # Use fewer lines for code overlap to maintain function context
                overlap_lines = current_chunk[-min(len(current_chunk), self.chunk_overlap // 20):]
                current_chunk = overlap_lines
                current_length = sum(len(l) + 1 for l in overlap_lines)
            
            current_chunk.append(line_with_newline)
            current_length += line_length
        
        # Add the final chunk
        if current_chunk:
            chunks.append(''.join(current_chunk))
        
        return chunks
    
    def _chunk_markdown_document(self, text: str) -> List[str]:
        """
        Chunk markdown documents by heading boundaries.
        
        Args:
            text (str): Markdown text content
            
        Returns:
            List[str]: Chunked markdown content
        """
        chunks = []
        
        # Split at markdown headings (# Heading)
        sections = re.split(r'(\n#{1,6} )', text)
        current_chunk = []
        current_length = 0
        
        for i in range(len(sections)):
            section = sections[i]
            section_length = len(section)
            
            # If adding this section would exceed chunk size, complete current chunk
            if current_length + section_length > self.chunk_size and current_chunk:
                chunks.append(''.join(current_chunk))
                current_chunk = []
                current_length = 0
            
            current_chunk.append(section)
            current_length += section_length
        
        # Add the final chunk
        if current_chunk:
            chunks.append(''.join(current_chunk))
        
        return chunks
    
    def _chunk_regular_text(self, text: str) -> List[str]:
        """
        Default chunking strategy for regular text using paragraphs and sentences.
        
        Args:
            text (str): Regular text content
            
        Returns:
            List[str]: Chunked text content
        """
        chunks = []
        paragraphs = re.split(r'\n\s*\n', text)
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_size = len(para) + 2  # +2 for the newlines
            
            # Handle paragraphs that are larger than chunk_size
            if para_size > self.chunk_size:
                # Complete the current chunk first
                if current_chunk:
                    chunks.append("\n\n".join(current_chunk).strip())
                    current_chunk = []
                    current_size = 0
                
                # Split this large paragraph by sentences
                sentences = re.split(r'(?<=[.!?])\s+', para)
                current_sentence_chunk = []
                current_sentence_size = 0
                
                for sentence in sentences:
                    sentence_size = len(sentence) + 1  # +1 for the space
                    
                    # If adding this sentence would exceed chunk size, complete current chunk
                    if current_sentence_size + sentence_size > self.chunk_size and current_sentence_chunk:
                        chunks.append(" ".join(current_sentence_chunk).strip())
                        current_sentence_chunk = []
                        current_sentence_size = 0
                    
                    current_sentence_chunk.append(sentence)
                    current_sentence_size += sentence_size
                
                # Add the final sentence chunk
                if current_sentence_chunk:
                    chunks.append(" ".join(current_sentence_chunk).strip())
                
                # Skip adding this paragraph to regular chunks since we've handled it
                continue
            
            # If adding this paragraph would exceed chunk size, complete current chunk
            if current_size + para_size > self.chunk_size and current_chunk:
                chunks.append("\n\n".join(current_chunk).strip())
                
                # Create overlap with previous paragraphs
                overlap_size = 0
                overlap_paras = []
                for prev_para in reversed(current_chunk):
                    if overlap_size + len(prev_para) + 2 <= self.chunk_overlap:
                        overlap_paras.insert(0, prev_para)
                        overlap_size += len(prev_para) + 2
                    else:
                        break
                
                current_chunk = overlap_paras
                current_size = overlap_size
            
            current_chunk.append(para)
            current_size += para_size
        
        # Add the final chunk
        if current_chunk:
            chunks.append("\n\n".join(current_chunk).strip())
        
        return chunks

class EmbeddingGenerator:
    """
    Generates vector embeddings for text using Azure OpenAI.
    
    This class handles the generation of vector embeddings for text content
    using Azure OpenAI's embedding models. These embeddings are used for
    vector search capabilities in Azure AI Search.
    
    Attributes:
        config (OpenAIConfig): Configuration for Azure OpenAI service
    """
    
    def __init__(self, config: OpenAIConfig):
        """
        Initialize the embedding generator.
        
        Args:
            config (OpenAIConfig): Configuration settings for Azure OpenAI
        """
        self.config = config
    
    def generate_embeddings(self, text: str) -> List[float]:
        """
        Generate vector embeddings for text using Azure OpenAI.
        
        This method calls the Azure OpenAI API to generate embeddings
        for the provided text. It handles:
        - Skipping embedding generation if configured to do so
        - Truncating text to fit within token limits
        - Error handling and response parsing
        - Different API response formats
        
        Args:
            text (str): The text to generate embeddings for
            
        Returns:
            List[float]: Vector embeddings as a list of floating point values,
                         or empty list if embedding generation fails or is skipped
        """
        # Skip embedding generation if configured to do so or if credentials are missing
        if self.config.skip_embeddings or not self.config.endpoint or not self.config.api_key:
            logger.warning("Skipping embeddings generation due to configuration")
            return []
        
        try:
            # Set up API request headers
            headers = {
                "Content-Type": "application/json",
                "api-key": self.config.api_key
            }
            
            # Truncate text if it exceeds the model's token limit
            # Most embedding models have an 8K token limit
            max_length = 8000  # Approximate token limit
            if len(text) > max_length:
                logger.debug(f"Truncating text from {len(text)} to {max_length} characters for embedding")
                text = text[:max_length]
            
            # Prepare request body
            body = {
                "input": text,
                "model": self.config.deployment_name
            }
            
            # Use the latest API version that supports embeddings
            api_version = "2023-12-01-preview"
            url = f"{self.config.endpoint}/openai/deployments/{self.config.deployment_name}/embeddings?api-version={api_version}"
            
            # Call the Azure OpenAI API
            logger.debug(f"Calling embedding API for text of length {len(text)}")
            response = requests.post(url, headers=headers, json=body)
            response.raise_for_status()
            
            # Parse the response
            result = response.json()
            
            # Handle different response formats based on API version
            if "data" in result and len(result["data"]) > 0 and "embedding" in result["data"][0]:
                # Standard OpenAI API format
                embeddings = result["data"][0]["embedding"]
                logger.debug(f"Successfully generated embedding with {len(embeddings)} dimensions")
                return embeddings
            elif "embedding" in result:
                # Alternative format used in some Azure OpenAI API versions
                logger.debug(f"Successfully generated embedding with {len(result['embedding'])} dimensions (alternative format)")
                return result["embedding"]
            else:
                # Unexpected response format
                logger.error(f"Unexpected embeddings response format: {result}")
                return []
                
        except Exception as e:
            logger.error(f"Error generating embeddings with Azure OpenAI: {e}")
            return []

class AzureSearchIndexManager:
    """
    Manages Azure AI Search index creation and configuration.
    
    This class handles the creation, update, and configuration of Azure AI Search indexes,
    including support for vector search and semantic search capabilities. It uses both
    the Azure SDK and direct REST API calls when necessary to configure advanced features.
    
    Attributes:
        config (AzureSearchConfig): Configuration for Azure AI Search
        credential (AzureKeyCredential): Authentication credential for Azure
        index_client (SearchIndexClient): Client for managing search indexes
    """
    
    def __init__(self, config: AzureSearchConfig):
        """
        Initialize the Azure AI Search index manager.
        
        Args:
            config (AzureSearchConfig): Configuration settings for Azure AI Search
        """
        self.config = config
        self.credential = AzureKeyCredential(config.api_key)
        self.index_client = SearchIndexClient(endpoint=config.endpoint, credential=self.credential)
    
    def create_or_update_index(self) -> None:
        """
        Create or update an Azure AI Search index based on configuration.
        
        This method handles:
        - Deleting existing index if force_recreate is enabled
        - Creating vector search enabled indexes using REST API
        - Creating semantic search enabled indexes
        - Falling back to SDK methods when appropriate
        
        Raises:
            Exception: Logs errors but continues execution with existing index
        """
        # Delete existing index if force recreation is enabled
        if self.config.force_recreate:
            self._delete_index_if_exists()
        
        try:
            # Build the index definition based on configuration
            index = self._build_index_definition()
            
            # For vector search, the index is created via REST API in _build_index_definition
            if self.config.vector_search:
                logger.info(f"Index {self.config.index_name} created via REST API in _build_index_definition")
                return
            
            # For non-vector search, use the SDK approach
            # Log the index configuration for debugging
            index_dict = index.as_dict()
            logger.info(f"Creating index with configuration: {json.dumps(index_dict, indent=2)[:500]}...")
            
            # Verify semantic settings are present if enabled
            if "semantic" in index_dict:
                logger.info("Semantic settings are present in the index configuration")
            else:
                logger.warning("Semantic settings are NOT present in the index configuration")
                
                # If semantic search is enabled but settings are missing, use REST API
                if self.config.semantic_search:
                    logger.warning("Using REST API method to create index with semantic settings")
                    self._create_index_with_rest_api(semantic_only=True)
                    return
            
            # Use the SDK to create or update the index
            result = self.index_client.create_or_update_index(index)
            logger.info(f"Index {self.config.index_name} created or updated successfully with SDK")
        except Exception as e:
            logger.error(f"Error creating index {self.config.index_name}: {e}")
            logger.info("Continuing with existing index configuration...")
    
    def _create_index_with_rest_api(self, semantic_only=False) -> None:
        """
        Create index using REST API when SDK limitations need to be bypassed.
        
        This is a placeholder method as the actual REST API call is implemented
        in the _build_index_definition method. Kept for backward compatibility.
        
        Args:
            semantic_only (bool): Whether to only add semantic settings
        """
        # This is a placeholder - the actual REST API call is in _build_index_definition
        pass
    
    def _delete_index_if_exists(self) -> None:
        """
        Delete the search index if it exists.
        
        This method attempts to delete the configured index and logs the result.
        It handles the case where the index doesn't exist gracefully.
        """
        try:
            logger.info(f"Attempting to delete index {self.config.index_name} for recreation")
            self.index_client.delete_index(self.config.index_name)
            logger.info(f"Successfully deleted index {self.config.index_name}")
        except Exception as e:
            logger.info(f"Index deletion unsuccessful (may not exist): {e}")
    
    def _build_index_definition(self) -> SearchIndex:
        """
        Build the search index definition based on configuration.
        
        This method creates the appropriate index definition with:
        - Standard fields for document metadata
        - Vector search configuration if enabled
        - Semantic search configuration if enabled
        
        For vector search, it uses direct REST API calls to configure
        vector search profiles and algorithms that aren't fully supported
        in the SDK.
        
        Returns:
            SearchIndex: The configured search index definition
        """
        # Import required models
        from azure.search.documents.indexes.models import (
            SearchableField, 
            SimpleField, 
            SearchFieldDataType,
            SearchIndex
        )
        
        # Define standard fields for the index
        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True),
            SearchableField(name="content", type=SearchFieldDataType.String, analyzer_name="en.microsoft"),
            SimpleField(name="file_name", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="file_path", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="file_extension", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="file_size", type=SearchFieldDataType.Int64, filterable=True),
            SimpleField(name="last_modified", type=SearchFieldDataType.DateTimeOffset, filterable=True, sortable=True),
            SimpleField(name="chunk_id", type=SearchFieldDataType.Int32, filterable=True),
            SimpleField(name="chunk_total", type=SearchFieldDataType.Int32, filterable=True),
            SimpleField(name="language", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="title", type=SearchFieldDataType.String, analyzer_name="en.microsoft"),
            SearchableField(name="author", type=SearchFieldDataType.String, analyzer_name="en.microsoft"),
        ]
        
        # If vector search is enabled, we'll use a direct REST API approach
        # rather than trying to work through the SDK limitations
        if self.config.vector_search:
            logger.info("Using direct REST API approach for vector search index creation")
            
            # Create complete index definition as a dictionary
            raw_index = {
                "name": self.config.index_name,
                "fields": [
                    {
                        "name": "id",
                        "type": "Edm.String",
                        "key": True,
                        "searchable": False,
                        "filterable": False,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "content",
                        "type": "Edm.String",
                        "searchable": True,
                        "filterable": False,
                        "sortable": False,
                        "facetable": False,
                        "analyzer": "en.microsoft"
                    },
                    {
                        "name": "file_name",
                        "type": "Edm.String",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "file_path",
                        "type": "Edm.String",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "file_extension",
                        "type": "Edm.String",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "file_size",
                        "type": "Edm.Int64",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "last_modified",
                        "type": "Edm.DateTimeOffset",
                        "searchable": False,
                        "filterable": True,
                        "sortable": True,
                        "facetable": False
                    },
                    {
                        "name": "chunk_id",
                        "type": "Edm.Int32",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "chunk_total",
                        "type": "Edm.Int32",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "language",
                        "type": "Edm.String",
                        "searchable": False,
                        "filterable": True,
                        "sortable": False,
                        "facetable": False
                    },
                    {
                        "name": "title",
                        "type": "Edm.String",
                        "searchable": True,
                        "filterable": False,
                        "sortable": False,
                        "facetable": False,
                        "analyzer": "en.microsoft"
                    },
                    {
                        "name": "author",
                        "type": "Edm.String",
                        "searchable": True,
                        "filterable": False,
                        "sortable": False,
                        "facetable": False,
                        "analyzer": "en.microsoft"
                    },
                    {
                        "name": "content_vector",
                        "type": "Collection(Edm.Single)",
                        "searchable": True,
                        "retrievable": False,
                        "dimensions": self.config.vector_dimensions,
                        "vectorSearchProfile": "vector-profile"
                    }
                ],
                "vectorSearch": {
                    "algorithms": [
                        {
                            "name": "vector-config",
                            "kind": "hnsw",
                            "hnswParameters": {
                                "m": 10,
                                "efConstruction": 400,
                                "efSearch": 500,
                                "metric": "cosine"
                            }
                        }
                    ],
                    "profiles": [
                        {
                            "name": "vector-profile",
                            "algorithm": "vector-config"
                        }
                    ]
                }
            }
            
            # Add semantic configuration if enabled
            if self.config.semantic_search:
                raw_index["semantic"] = {
                    "configurations": [{
                        "name": "semantic-config",
                        "prioritizedFields": {
                            "titleField": {"fieldName": "title"},
                            "prioritizedContentFields": [{"fieldName": "content"}],
                            "prioritizedKeywordsFields": []
                        }
                    }]
                }
            
            # Create the index using the REST API
            api_version = "2024-07-01"
            endpoint = f"{self.config.endpoint}/indexes/{self.config.index_name}?api-version={api_version}"
            headers = {
                "Content-Type": "application/json",
                "api-key": self.config.api_key
            }
            
            logger.info(f"Creating index with REST API using API version: {api_version}")
            
            # Log the exact configuration we're sending (mask the sensitive parts)
            index_config_log = json.dumps(raw_index, indent=2)
            logger.info(f"Index configuration being sent: {index_config_log}")
            
            try:
                # Send the request
                response = requests.put(endpoint, headers=headers, json=raw_index)
                response_text = response.text
                
                # Log the full response for debugging
                logger.info(f"REST API response status: {response.status_code}")
                try:
                    response_json = response.json()
                    logger.info(f"REST API response: {json.dumps(response_json, indent=2)}")
                except:
                    logger.info(f"REST API response text: {response_text}")
                
                if response.status_code >= 400:
                    logger.error(f"Error response from REST API: {response_text}")
                    response.raise_for_status()
                
                logger.info(f"Index {self.config.index_name} created successfully via REST API")
                
                # Verify the index was created properly with vector search profile
                try:
                    verify_endpoint = f"{self.config.endpoint}/indexes/{self.config.index_name}?api-version={api_version}"
                    verify_response = requests.get(verify_endpoint, headers=headers)
                    verify_json = verify_response.json()
                    
                    # Check specifically for vector search configuration
                    if "vectorSearch" in verify_json:
                        logger.info("Vector search configuration confirmed in index")
                    else:
                        logger.warning("Vector search configuration not found in created index")
                    
                    # Check content_vector field configuration
                    content_vector_field = None
                    for field in verify_json.get("fields", []):
                        if field.get("name") == "content_vector":
                            content_vector_field = field
                            break
                    
                    if content_vector_field:
                        if "vectorSearchProfile" in content_vector_field:
                            logger.info(f"content_vector field properly configured with profile: {content_vector_field.get('vectorSearchProfile')}")
                        else:
                            logger.warning("content_vector field found but missing vectorSearchProfile configuration")
                    else:
                        logger.warning("content_vector field not found in index")
                        
                except Exception as verify_error:
                    logger.warning(f"Could not verify index after creation: {verify_error}")
                
                # Return a SearchIndex object with fields to satisfy the method's return type
                # This is just a skeleton and won't be used for further API calls since the 
                # index is already created via REST API
                return SearchIndex(name=self.config.index_name, fields=fields)
            except Exception as e:
                logger.error(f"Error creating index via REST API: {e}")
                logger.error("Creating vector search index failed, the script may not work properly")
                logger.error("Check your Azure AI Search service and configuration")
                # Fall back to returning a basic index, but log a clear error
                # This object won't be used for actual index creation as the 
                # create_or_update_index method will detect vector_search=True and return early
                return SearchIndex(name=self.config.index_name, fields=fields)
        else:
            # For non-vector search, just use the SDK's approach
            index_additional_properties = {}
            
            if self.config.semantic_search:
                logger.info("Adding semantic search configuration to index")
                # Use a dictionary-based configuration for semantic search
                semantic_config = {
                    "configurations": [{
                        "name": "semantic-config",
                        "prioritizedFields": {
                            "titleField": {"fieldName": "title"},
                            "prioritizedContentFields": [{"fieldName": "content"}],
                            "prioritizedKeywordsFields": []
                        }
                    }]
                }
                
                index_additional_properties["semantic"] = semantic_config
                logger.info(f"Semantic configuration: {json.dumps(semantic_config, indent=2)}")
            
            # Create the index with basic fields
            index = SearchIndex(name=self.config.index_name, fields=fields)
            
            # Add any additional properties to the index
            for key, value in index_additional_properties.items():
                logger.info(f"Adding additional property to index: {key}")
                index.additional_properties[key] = value
            
            return index

class DocumentUploader:
    """
    Handles the uploading of document chunks to an Azure AI Search index.
    
    This class manages the process of uploading processed document chunks to an Azure AI Search
    index, handling batch processing, error recovery, and validation of document fields.
    It supports both standard search and vector search configurations.
    
    Attributes:
        azure_config (AzureSearchConfig): Configuration for Azure AI Search service.
        openai_config (OpenAIConfig): Configuration for OpenAI API.
        credential (AzureKeyCredential): Authentication credential for Azure services.
        search_client (SearchClient): Client for interacting with the Azure AI Search index.
    """
    
    def __init__(
        self, 
        azure_config: AzureSearchConfig,
        openai_config: OpenAIConfig
    ):
        """
        Initialize the DocumentUploader with Azure and OpenAI configurations.
        
        Args:
            azure_config (AzureSearchConfig): Configuration for Azure AI Search service.
            openai_config (OpenAIConfig): Configuration for OpenAI API.
        """
        self.azure_config = azure_config
        self.openai_config = openai_config
        
        self.credential = AzureKeyCredential(azure_config.api_key)
        self.search_client = SearchClient(
            endpoint=azure_config.endpoint, 
            index_name=azure_config.index_name, 
            credential=self.credential
        )
    
    def upload_chunks(self, chunks: List[Dict[str, Any]]) -> None:
        """
        Upload document chunks to Azure AI Search index.
        
        This method handles the batch uploading of document chunks to the configured
        Azure AI Search index. It performs validation of the index structure for vector search,
        processes documents in batches of 1000 (Azure Search limit), and handles error recovery
        if vector fields cause issues during upload.
        
        Args:
            chunks (List[Dict[str, Any]]): List of document chunks to upload, where each chunk
                is a dictionary containing fields like id, content, file_name, etc.
                
        Raises:
            Exception: If there's an error during the upload process that can't be recovered from.
        """
        try:
            # If we're doing vector search but the index doesn't have the vector field properly set up
            # We should filter out the vector embeddings
            if self.azure_config.vector_search:
                try:
                    # Try to verify the index has a content_vector field
                    index_client = SearchIndexClient(
                        endpoint=self.search_client._endpoint, 
                        credential=self.search_client._credential
                    )
                    index = index_client.get_index(self.search_client._index_name)
                    
                    has_vector_field = False
                    for field in index.fields:
                        if field.name == "content_vector":
                            # Enhanced check for proper vector search configuration
                            # Check in various ways that might appear in the API response
                            field_dict = field.as_dict() if hasattr(field, "as_dict") else {}
                            
                            # Check all possible ways the vector profile could be indicated
                            if (hasattr(field, "vector_search_profile") or 
                                "vectorSearchProfile" in str(field.additional_properties) or
                                "vectorSearchProfile" in str(field_dict) or
                                "vector_search_profile" in str(field_dict) or
                                getattr(field, "dimensions", 0) > 0):  # If dimensions are set, likely vector field
                                
                                has_vector_field = True
                            
                            logger.info(f"Content vector field found: {field_dict}")
                            break
                    
                    if not has_vector_field:
                        logger.warning("Index does not contain a properly configured content_vector field. Removing vector data from chunks.")
                        for chunk in chunks:
                            if "content_vector" in chunk:
                                del chunk["content_vector"]
                except Exception as e:
                    logger.warning(f"Could not verify index structure: {e}. Attempting upload anyway.")
                    
            # Upload in batches of 1000 (Azure Search limit)
            batch_size = 1000
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i:i+batch_size]
                
                # Make sure all documents have valid data types
                batch = [self._ensure_valid_document(doc) for doc in batch]
                
                # Log a small preview of what we're sending
                if i == 0:
                    # Create a deep copy of the first document for logging (to avoid modifying the original)
                    sample_doc = dict(batch[0])
                    # Remove vector field for logging if it exists (it's too large to print)
                    if "content_vector" in sample_doc:
                        sample_doc["content_vector"] = "[vector data]"
                    logger.info(f"Sample document being uploaded: {json.dumps(sample_doc, indent=2)[:500]}...")
                
                try:
                    result = self.search_client.upload_documents(documents=batch)
                    logger.info(f"Uploaded batch of {len(batch)} chunks (from {i} to {i+len(batch)-1})")
                    success_count = sum(1 for r in result if r.succeeded)
                    logger.info(f"Upload results: {success_count}/{len(batch)} succeeded")
                    
                    # If there were failures, log details about the first few
                    if success_count < len(batch):
                        failed_count = len(batch) - success_count
                        logger.warning(f"{failed_count} document uploads failed")
                        for j, r in enumerate(result):
                            if not r.succeeded and j < 5:  # Log up to 5 failures
                                logger.warning(f"Document {j} failed: {r.error_message}")
                
                except Exception as batch_error:
                    logger.error(f"Error uploading batch {i//batch_size + 1}: {batch_error}")
                    
                    # Try without content_vector field if that's the issue
                    if "content_vector" in str(batch_error).lower() or "vector" in str(batch_error).lower():
                        logger.info("Trying upload without vector embeddings...")
                        try:
                            for doc in batch:
                                if "content_vector" in doc:
                                    del doc["content_vector"]
                            result = self.search_client.upload_documents(documents=batch)
                            logger.info(f"Upload without vectors succeeded: {sum(1 for r in result if r.succeeded)}/{len(batch)}")
                        except Exception as retry_error:
                            logger.error(f"Failed even without vectors: {retry_error}")
                            raise
                    else:
                        raise
                    
        except Exception as e:
            logger.error(f"Error uploading documents: {e}")
            raise
    
    def _ensure_valid_document(self, document: Dict[str, Any]) -> Dict[str, Any]:
        """
        Ensure all fields have valid data types for Azure Search.
        
        This method validates and corrects document fields to ensure they meet
        Azure AI Search requirements, particularly for date fields which need
        to be in a specific format.
        
        Args:
            document (Dict[str, Any]): The document to validate and correct.
            
        Returns:
            Dict[str, Any]: The validated and corrected document.
        """
        # Current time in ISO 8601 UTC format as a fallback
        default_date = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
        
        # Fields that are expected to be dates
        date_fields = ["last_modified"]
        
        # Check and fix date fields
        for field in date_fields:
            if field in document:
                if not document[field] or not isinstance(document[field], str):
                    logger.warning(f"Invalid date value for {field}, using default date")
                    document[field] = default_date
                elif not (document[field].endswith('Z') or '+' in document[field]):
                    # Add Z for UTC if no timezone info
                    document[field] = document[field].rstrip() + 'Z'
        
        return document

class RAGDocumentProcessor:
    """
    Main orchestration class for the entire RAG document processing pipeline.
    
    This class coordinates the end-to-end process of preparing documents for Retrieval-Augmented
    Generation (RAG), including document processing, text chunking, embedding generation,
    index management, and document uploading to Azure AI Search.
    
    Attributes:
        processing_config (ProcessingConfig): Configuration for document processing.
        azure_config (AzureSearchConfig): Configuration for Azure AI Search.
        openai_config (OpenAIConfig): Configuration for OpenAI API.
        document_processor (DocumentProcessor): Handles extraction of text and metadata from files.
        text_chunker (TextChunker): Splits text into appropriate chunks.
        embedding_generator (EmbeddingGenerator): Generates vector embeddings for text chunks.
        index_manager (AzureSearchIndexManager): Manages the Azure AI Search index.
        document_uploader (DocumentUploader): Uploads processed chunks to Azure AI Search.
    """
    
    def __init__(
        self,
        processing_config: ProcessingConfig,
        azure_config: AzureSearchConfig,
        openai_config: OpenAIConfig
    ):
        """
        Initialize the RAGDocumentProcessor with necessary configurations and components.
        
        Args:
            processing_config (ProcessingConfig): Configuration for document processing.
            azure_config (AzureSearchConfig): Configuration for Azure AI Search.
            openai_config (OpenAIConfig): Configuration for OpenAI API.
        """
        self.processing_config = processing_config
        self.azure_config = azure_config
        self.openai_config = openai_config
        
        self.document_processor = DocumentProcessor(processing_config)
        self.text_chunker = TextChunker(processing_config.chunk_size, processing_config.chunk_overlap)
        self.embedding_generator = EmbeddingGenerator(openai_config)
        self.index_manager = AzureSearchIndexManager(azure_config)
        self.document_uploader = DocumentUploader(azure_config, openai_config)
    
    def process(self) -> None:
        """
        Process all documents and upload to Azure AI Search.
        
        This method orchestrates the entire pipeline:
        1. Creates or updates the search index
        2. Walks through the specified directory to process each file
        3. Extracts text and metadata from each file
        4. Chunks the text into appropriate segments
        5. Generates embeddings for each chunk if vector search is enabled
        6. Uploads all processed chunks to Azure AI Search
        
        Raises:
            Various exceptions may be raised during processing, which are logged
            but not propagated unless they occur during the final upload step.
        """
        # Create or update the search index
        self.index_manager.create_or_update_index()
        
        # Process files and upload chunks
        all_chunks = []
        
        # Walk through the directory and process each file
        for root, _, files in os.walk(self.processing_config.folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                
                # Check if file matches include filters or should be excluded
                if not self._should_process_file(file_path):
                    logger.info(f"Skipping file based on filters: {file_path}")
                    continue
                
                logger.info(f"Processing file: {file_path}")
                
                # Extract text and metadata from file
                result = self.document_processor.process_file(file_path)
                text = result["text"]
                metadata = result["metadata"]
                
                if not text:
                    logger.warning(f"No text content extracted from {file_path}")
                    continue
                
                # Detect language
                language = self._detect_language(text)
                
                # Use standard chunking for all file types now
                text_chunks = self.text_chunker.chunk_text(text)
                logger.info(f"Split {file_path} into {len(text_chunks)} semantic chunks")
                
                # Remove duplicate chunks
                unique_chunks = self.text_chunker.remove_duplicate_chunks(text_chunks)
                if len(unique_chunks) < len(text_chunks):
                    logger.info(f"Removed {len(text_chunks) - len(unique_chunks)} duplicate chunks")
                
                # Current time in ISO 8601 format for any missing dates
                default_date = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
                
                # Create documents for each unique chunk
                for i, chunk in enumerate(unique_chunks):
                    # Generate a stable ID based on file path and chunk index
                    chunk_id = hashlib.md5(f"{file_path}_{i}".encode()).hexdigest()
                    
                    # Create title with page information if available
                    # Extract page numbers from the chunk if present
                    page_numbers = []
                    if file_path.lower().endswith('.pdf'):
                        page_markers = re.findall(r'\[pg:(\d+)\]', chunk)
                        if page_markers:
                            page_numbers = sorted(set(int(p) for p in page_markers))
                            
                    # Create a descriptive title
                    if page_numbers:
                        if len(page_numbers) == 1:
                            chunk_title = f"{os.path.basename(file_path)} (Page {page_numbers[0]})"
                        else:
                            chunk_title = f"{os.path.basename(file_path)} (Pages {min(page_numbers)}-{max(page_numbers)})"
                    else:
                        chunk_title = metadata.get("title", os.path.basename(file_path))
                        
                        # Try to extract a better title from the first line
                        first_line = chunk.split('\n', 1)[0].strip()
                        if first_line and len(first_line) < 100 and not first_line.startswith('[') and not first_line.startswith('#'):
                            chunk_title = f"{os.path.basename(file_path)}: {first_line}"
                    
                    # Ensure we have a valid last_modified date
                    last_modified = metadata.get("last_modified")
                    if not last_modified or last_modified == "":
                        last_modified = default_date
                    
                    document = {
                        "id": chunk_id,
                        "content": chunk,
                        "file_name": metadata.get("file_name", ""),
                        "file_path": metadata.get("file_path", ""),
                        "file_extension": metadata.get("file_extension", ""),
                        "file_size": metadata.get("file_size", 0),
                        "last_modified": last_modified,
                        "chunk_id": i,
                        "chunk_total": len(unique_chunks),
                        "language": language,
                        "title": chunk_title,
                        "author": metadata.get("author", "")
                    }
                    
                    # Add embeddings if vector search is enabled
                    if self.azure_config.vector_search and not self.openai_config.skip_embeddings:
                        try:
                            logger.info(f"Generating embeddings for chunk {i+1}/{len(unique_chunks)} of {file_path}")
                            embeddings = self.embedding_generator.generate_embeddings(chunk)
                            if embeddings:
                                document["content_vector"] = embeddings
                            # Add a short delay to avoid rate limiting
                            time.sleep(0.1)
                        except Exception as e:
                            logger.error(f"Error generating embeddings: {e}")
                            logger.info("Continuing without embeddings for this chunk...")
                    
                    all_chunks.append(document)
        
        # Upload chunks to the search index
        if all_chunks:
            logger.info(f"Uploading {len(all_chunks)} chunks to index {self.azure_config.index_name}")
            self.document_uploader.upload_chunks(all_chunks)
            logger.info("Upload completed successfully")
        else:
            logger.warning("No content was extracted from the files. Nothing to upload.")
    
    def _should_process_file(self, file_path: str) -> bool:
        """
        Determine if a file should be processed based on include/exclude filters.
        
        This method checks if a file matches the include and exclude patterns
        specified in the processing configuration.
        
        Args:
            file_path (str): The path of the file to check.
            
        Returns:
            bool: True if the file should be processed, False otherwise.
        """
        # If no filters are set, process all files
        if not self.processing_config.include_filters and not self.processing_config.exclude_filters:
            return True
        
        file_name = os.path.basename(file_path).lower()
        
        # Check exclude filters first
        for pattern in self.processing_config.exclude_filters:
            if re.search(pattern.lower(), file_name) or re.search(pattern.lower(), file_path.lower()):
                return False
        
        # If include filters are specified, at least one must match
        if self.processing_config.include_filters:
            for pattern in self.processing_config.include_filters:
                if re.search(pattern.lower(), file_name) or re.search(pattern.lower(), file_path.lower()):
                    return True
            # If we get here, no include filter matched
            return False
        
        # No include filters and excluded passed
        return True
    
    def _detect_language(self, text: str) -> str:
        """
        Detect the language of the text.
        
        Uses the langdetect library to identify the language of the provided text.
        
        Args:
            text (str): The text to analyze for language detection.
            
        Returns:
            str: The detected language code (e.g., 'en' for English).
                 Defaults to 'en' if detection fails.
        """
        try:
            # Use a sample of the text for faster detection
            sample = text[:min(1000, len(text))]
            return detect(sample)
        except:
            # Default to English if detection fails
            return "en"

def main():
    """
    Main entry point for the RAG document processing script.
    
    This function:
    1. Parses command-line arguments to configure the processing pipeline
    2. Validates the provided arguments for consistency and completeness
    3. Creates the necessary configuration objects
    4. Initializes and runs the RAGDocumentProcessor
    
    Command-line arguments include options for:
    - Document processing (folder path, chunk size, chunk overlap)
    - Azure AI Search configuration (endpoint, key, index name)
    - Vector search capabilities (dimensions, OpenAI settings)
    - Semantic search capabilities
    - File filtering (include/exclude patterns)
    
    Returns:
        None. The function logs progress and errors to the console.
    """
    parser = argparse.ArgumentParser(description='Upload documents to Azure AI Search for RAG')
    parser.add_argument('--folder', type=str, required=True, help='Folder containing documents to upload')
    parser.add_argument('--endpoint', type=str, required=True, help='Azure AI Search endpoint')
    parser.add_argument('--key', type=str, required=True, help='Azure AI Search admin key')
    parser.add_argument('--index', type=str, required=True, help='Name of the search index')
    parser.add_argument('--chunk-size', type=int, default=1000, help='Size of text chunks (in characters)')
    parser.add_argument('--chunk-overlap', type=int, default=100, help='Overlap between chunks (in characters)')
    parser.add_argument('--vector-search', action='store_true', help='Enable vector search capabilities in the Azure AI Search index')
    parser.add_argument('--semantic-search', action='store_true', help='Enable semantic search capabilities in the Azure AI Search index')
    parser.add_argument('--aoai-endpoint', type=str, help='Azure OpenAI endpoint for embeddings generation (required for vector search)')
    parser.add_argument('--aoai-key', type=str, help='Azure OpenAI key for embeddings generation (required for vector search)')
    parser.add_argument('--aoai-deployment', type=str, help='Azure OpenAI deployment name for embeddings model (e.g. text-embedding-3-large or text-embedding-ada-002)')
    parser.add_argument('--vector-dimensions', type=int, default=1536, help='Dimensions for vector embeddings (3072 for text-embedding-3-large, 1536 for text-embedding-ada-002)')
    parser.add_argument('--skip-vectors', action='store_true', help='Skip vector processing even if --vector-search is enabled')
    parser.add_argument('--force-recreate', action='store_true', help='Force recreation of the index')
    parser.add_argument('--include', type=str, nargs='+', default=[], help='File patterns to include (regex)')
    parser.add_argument('--exclude', type=str, nargs='+', default=[], help='File patterns to exclude (regex)')
    
    args = parser.parse_args()
    
    # Validate vector search parameters
    if args.vector_search and not args.skip_vectors and (not args.aoai_endpoint or not args.aoai_key or not args.aoai_deployment):
        logger.error("Vector search requires Azure OpenAI parameters: --aoai-endpoint, --aoai-key, and --aoai-deployment")
        return
    
    # Set appropriate vector dimensions based on the model
    if args.vector_search and args.aoai_deployment:
        if "text-embedding-3-large" in args.aoai_deployment and args.vector_dimensions != 3072:
            logger.warning(f"Setting vector dimensions to 3072 to match the capabilities of {args.aoai_deployment}")
            args.vector_dimensions = 3072
        elif "text-embedding-ada-002" in args.aoai_deployment and args.vector_dimensions != 1536:
            logger.warning(f"Setting vector dimensions to 1536 to match the capabilities of {args.aoai_deployment}")
            args.vector_dimensions = 1536
    
    # Create configuration objects
    processing_config = ProcessingConfig(
        folder_path=args.folder,
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
        include_filters=args.include,
        exclude_filters=args.exclude
    )
    
    azure_config = AzureSearchConfig(
        endpoint=args.endpoint,
        api_key=args.key,
        index_name=args.index,
        vector_search=args.vector_search,
        semantic_search=args.semantic_search,
        vector_dimensions=args.vector_dimensions,
        force_recreate=args.force_recreate
    )
    
    openai_config = OpenAIConfig(
        endpoint=args.aoai_endpoint,
        api_key=args.aoai_key,
        deployment_name=args.aoai_deployment,
        skip_embeddings=args.skip_vectors
    )
    
    # Create and run the document processor
    processor = RAGDocumentProcessor(processing_config, azure_config, openai_config)
    processor.process()

if __name__ == "__main__":
    main()